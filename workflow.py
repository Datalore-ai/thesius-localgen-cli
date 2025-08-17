import json
import base64
import io
import pymupdf
import fitz
import os
from pptx import Presentation
from PIL import Image
from docx import Document
from mistralai import Mistral
import pdfplumber
import asyncio

from qdrant_setup import *
from agents.generation_agent import generation_agent
from agents.schema_agent import generate_dataset_schema
from agents.evolution_agent.evolver import evolve_dataset
from utils import process_datagen_prompt


client = Mistral(api_key=os.getenv("MISTRAL_API_KEY"))

def encode_pdf(pdf_bytes: bytes):
    """Encode PDF bytes to a base64 string."""
    try:
        return base64.b64encode(pdf_bytes).decode("utf-8")
    except Exception as e:
        print(f"Error encoding PDF to base64: {e}")
        return None

def convert_to_pdf(file_bytes: bytes, filename: str):
    extension = filename.lower().split('.')[-1]

    if extension == "pdf":
        return file_bytes

    buffer = io.BytesIO()
    pdf = fitz.open()

    if extension in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
        img.save(buffer, format="PDF")
        return buffer.getvalue()

    elif extension in {"txt", "md"}:
        text = file_bytes.decode("utf-8", errors="ignore")
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension in {"doc", "docx"}:
        doc = Document(io.BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs])
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            page = pdf.new_page()
            page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    else:
        raise ValueError(f"Unsupported file type: {extension}")
    
def process_page(idx, ocr_response=None):
    try:
        if ocr_response and hasattr(ocr_response, 'pages') and idx < len(ocr_response.pages):
            return ocr_response.pages[idx].markdown
        else:
            return f"Error: Page {idx + 1} not available in OCR response"
    except Exception as e:
        return f"Error processing page {idx + 1}: {e}"

def extract_text_from_pdf(pdf_bytes: bytes):
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
    except Exception:
        with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
            total_pages = len(doc)

    pdf_bytes = encode_pdf(pdf_bytes)

    try:
        response = client.ocr.process(
            model="mistral-ocr-latest",
            document={
                "type": "document_url",
                "document_url": f"data:application/pdf;base64,{pdf_bytes}"
            },
            include_image_base64=True
        )
    except Exception as e:
        return [f"Error during OCR processing: {e}"]

    page_numbers = list(range(total_pages))

    extracted_text = []
    for idx in page_numbers:
        extracted_text.append(process_page(idx, ocr_response=response))
        
    return extracted_text

def create_chunks(directory_path: str):
    file_paths = [
        os.path.abspath(os.path.join(directory_path, f))
        for f in os.listdir(directory_path)
        if os.path.isfile(os.path.join(directory_path, f))
    ]
    Chunks = []
    for idx, file_path in enumerate(file_paths):
        filename = os.path.basename(file_path)
        with open(file_path, "rb") as f:
            file_bytes = f.read()
        
        print(f"📑 Converting {filename} to pdf if necessary...")
        converted_pdf_bytes = convert_to_pdf(file_bytes, filename)

        print(f"📑 Extracting pages from file {filename}...")
        pages = extract_text_from_pdf(converted_pdf_bytes)
        for page in pages:
            Chunks.append({"filename": filename, "page_number": idx + 1,"page_content": page})
    return Chunks

def create_records(page_data: str, system_prompt: str):
    try:
        datarecords = generation_agent(page_data, system_prompt=system_prompt)
        return datarecords
    except Exception as e:
        print(f"QA generation failed for a page: {str(e)}")
    return []

async def generate_full_dataset(directory_path: str, system_prompt: str):
    Chunks = create_chunks(directory_path)

    dataset = []
    
    yield f"⚙️ Setting things up...\n\n"
    rag_pipeline_setup(user_id="test_user", documents=Chunks)

    Temp_Chunks = Chunks.copy()
    while len(Temp_Chunks) != 0:
        print(f"🧠 Generating your dataset - {int((len(Chunks)-len(Temp_Chunks))/len(Chunks) * 100)} % done")
        idx, current_chunk = select_random_chunk(Temp_Chunks)
        results = retrieve_from_store(current_chunk, user_id="test_user")

        # Context prep
        context = "\n\n\n\n".join(f"filename:{result.payload['document']['filename']}\nPage_number:{result.payload['document']['page_number']}\nPage_Content: {result.payload['document']["page_content"]}" for result in results)

        page_qas = create_records(context, system_prompt)
        dataset.extend(page_qas)
        page_qas = evolve_dataset(page_qas)
        dataset.extend(page_qas)

        similar_chunks = [result.payload['document'] for result in results]

        for chunk in similar_chunks: 
            if chunk in Temp_Chunks:
                Temp_Chunks.remove(chunk)
    
    remove_data_from_store(user_id="test_user")
        
    yield f"Dataset generation completed with {len(dataset)} rows!\n\n"
    yield f"data:__DONE__:{json.dumps({'rows': dataset})}\n\n"
